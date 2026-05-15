import os
import io
import asyncio
from typing import Optional, Dict, Any, List, Tuple
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import pysrt
import webvtt
from app.app_core.exceptions import FileProcessingError
from app.app_services.file_service import FileService
import edge_tts


class TextConversionService:
    """Service for handling text extraction from various document formats."""
    
    # Supported input formats
    SUPPORTED_INPUT_FORMATS = {
        'DOCX', 'DOC', 'PPTX', 'PPT', 'PDF', 'SRT', 'VTT', 'TXT'
    }
    
    # Popular Edge TTS voices
    POPULAR_VOICES = {
        # English (US)
        "female_us_aria": "en-US-AriaNeural",
        "female_us_jenny": "en-US-JennyNeural",
        "female_us_emma": "en-US-EmmaNeural",
        "male_us_guy": "en-US-GuyNeural",
        "male_us_andrew": "en-US-AndrewNeural",
        "male_us_brian": "en-US-BrianNeural",
        "male_us_christopher": "en-US-ChristopherNeural",
        "male_us_eric": "en-US-EricNeural",
        
        # English (UK)
        "female_uk_sonia": "en-GB-SoniaNeural",
        "female_uk_libby": "en-GB-LibbyNeural",
        "male_uk_thomas": "en-GB-ThomasNeural",
        "male_uk_ryan": "en-GB-RyanNeural",
        
        # Urdu & Hindi (Expert)
        "urdu_male": "ur-PK-AsadNeural",
        "urdu_female": "ur-PK-UzmaNeural",
        "hindi_male": "hi-IN-MadhurNeural",
        "hindi_female": "hi-IN-SwaraNeural",
        
        # Arabic
        "arabic_male": "ar-SA-HamedNeural",
        "arabic_female": "ar-SA-ZariyahNeural",
        
        # Spanish
        "spanish_male": "es-ES-AlvaroNeural",
        "spanish_female": "es-ES-ElviraNeural",
        
        # French
        "french_male": "fr-FR-HenriNeural",
        "french_female": "fr-FR-DeniseNeural",
        
        # German
        "german_male": "de-DE-ConradNeural",
        "german_female": "de-DE-KatjaNeural",

        # Chinese
        "chinese_male": "zh-CN-YunxiNeural",
        "chinese_female": "zh-CN-XiaoxiaoNeural",
        
        # Aliases for convenience
        "female_us": "en-US-AriaNeural",
        "male_us": "en-US-GuyNeural",
        "female_uk": "en-GB-SoniaNeural",
        "male_uk": "en-GB-ThomasNeural"
    }
    
    @staticmethod
    def word_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from Word document."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input Word file not found: {input_path}")
            
            # Load Word document
            doc = Document(input_path)
            
            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Word to text conversion failed: {str(e)}")
    
    @staticmethod
    def powerpoint_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input PowerPoint file not found: {input_path}")
            
            # Load PowerPoint presentation
            prs = Presentation(input_path)
            
            # Extract text from all slides
            text_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.extend(slide_text)
                    text_content.append("")  # Add empty line between slides
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"PowerPoint to text conversion failed: {str(e)}")
    
    @staticmethod
    def pdf_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from PDF document."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input PDF file not found: {input_path}")
            
            # Open PDF document
            doc = fitz.open(input_path)
            
            # Extract text from all pages
            text_content = []
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                
                if page_text.strip():
                    text_content.append(f"--- Page {page_num + 1} ---")
                    text_content.append(page_text.strip())
                    text_content.append("")  # Add empty line between pages
            
            # Close document
            doc.close()
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"PDF to text conversion failed: {str(e)}")
    
    @staticmethod
    def srt_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from SRT subtitle file."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input SRT file not found: {input_path}")
            
            # Load SRT file
            subs = pysrt.open(input_path)
            
            # Extract text from all subtitles
            text_content = []
            for sub in subs:
                if sub.text.strip():
                    text_content.append(sub.text.strip())
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"SRT to text conversion failed: {str(e)}")
    
    @staticmethod
    def vtt_to_text(input_path: str, output_filename: Optional[str] = None) -> str:
        """Extract text from VTT subtitle file."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input VTT file not found: {input_path}")
            
            # Load VTT file
            vtt = webvtt.read(input_path)
            
            # Extract text from all captions
            text_content = []
            for caption in vtt:
                if caption.text.strip():
                    text_content.append(caption.text.strip())
            
            # Join all text content
            extracted_text = "\n".join(text_content)
            
            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".txt"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".txt")
            
            # Save text file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"VTT to text conversion failed: {str(e)}")
    
    @staticmethod
    async def text_to_speech(input_path: str, output_filename: Optional[str] = None, language: str = "en", voice: Optional[str] = None) -> str:
        """Convert document/text to speech (MP3) using edge-tts for high quality voices."""
        try:
            if not os.path.exists(input_path):
                raise FileProcessingError(f"Input file not found: {input_path}")
            
            ext = os.path.splitext(input_path)[1].lower()
            text = ""
            temp_txt_path = None
            
            # Extract text based on file extension
            if ext == '.txt':
                with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
            elif ext == '.docx':
                temp_txt_path = TextConversionService.word_to_text(input_path)
            elif ext == '.pdf':
                temp_txt_path = TextConversionService.pdf_to_text(input_path)
            elif ext in ['.pptx', '.ppt']:
                temp_txt_path = TextConversionService.powerpoint_to_text(input_path)
            elif ext == '.srt':
                temp_txt_path = TextConversionService.srt_to_text(input_path)
            elif ext == '.vtt':
                temp_txt_path = TextConversionService.vtt_to_text(input_path)
            else:
                # Try to read as plain text if unknown extension
                try:
                    with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                except Exception:
                    raise FileProcessingError(f"Unsupported file type for text-to-speech: {ext}")

            # If we used a converter, read the resulting text file
            if temp_txt_path:
                with open(temp_txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                # Cleanup temp text file
                if os.path.exists(temp_txt_path):
                    os.remove(temp_txt_path)

            if not text.strip():
                raise FileProcessingError("No text found in the input file to convert to speech.")

            # Generate output path
            if output_filename and output_filename.strip():
                output_path, _ = FileService.generate_output_path_with_filename(
                    output_filename.strip(), default_extension=".mp3"
                )
            else:
                output_path = FileService.get_output_path(input_path, ".mp3")
            
            # Determine voice
            # If voice is a key in POPULAR_VOICES, use the value
            # Otherwise, use the voice string directly or fallback to Aria
            selected_voice = TextConversionService.POPULAR_VOICES.get(voice, voice) or "en-US-AriaNeural"
            
            # Convert text to speech using edge-tts
            communicate = edge_tts.Communicate(text, selected_voice)
            await communicate.save(output_path)
            
            return output_path
            
        except Exception as e:
            # Cleanup temp text file on error if it exists
            if 'temp_txt_path' in locals() and temp_txt_path and os.path.exists(temp_txt_path):
                os.remove(temp_txt_path)
            
            if isinstance(e, FileProcessingError):
                raise e
            raise FileProcessingError(f"Text to speech conversion failed: {str(e)}")

    @staticmethod
    def get_supported_formats() -> List[str]:
        """Get list of supported input formats."""
        return list(TextConversionService.SUPPORTED_INPUT_FORMATS)
    
    @staticmethod
    def cleanup_temp_files(*file_paths: str) -> None:
        """Clean up temporary files."""
        for file_path in file_paths:
            FileService.cleanup_file(file_path)
